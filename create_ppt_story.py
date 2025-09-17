from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_slide(title, bullets, notes=None):
    layout = prs.slide_layouts[1]  # Title & Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(38)
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(22)
    if notes is not None:
        slide.notes_slide.notes_text_frame.clear()
        slide.notes_slide.notes_text_frame.text = notes

# Título
slide0 = prs.slides.add_slide(prs.slide_layouts[0])
slide0.shapes.title.text = "Cognitiva-AI — IA intermodal para cribado de Alzheimer"
slide0.placeholders[1].text = "MRI + clínico · calibración por cohorte · decisión coste-sensible (S2) · demo lista"
slide0.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
slide0.placeholders[1].text_frame.paragraphs[0].font.size = Pt(20)
slide0.notes_slide.notes_text_frame.text = "Presentación de 8 minutos. Resaltar problema social y promesa: sensibilidad con coste controlado."

# Slides (idénticos a StoryTelling_FINAL.md)
slides_payload = [
    ("El problema que no se ve",
     ["Diagnóstico tardío → necesidad de cribado sensible",
      "Optimización con coste clínico: FN:FP = 5:1",
      "Objetivo: alta sensibilidad sin disparar FP"],
     "Cribado no es diagnóstico definitivo; nuestra métrica operativa penaliza FN cinco veces más."),
    ("Datos y reto",
     ["OASIS-1 (transversal) · OASIS-2 (longitudinal, 1ª visita/paciente)",
      "MRI: 20 slices → 56 features por paciente",
      "Clínico: homogeneización + imputación + OHE(Sex) · anti-fuga (no CDR/Group como feature)"],
     "Estándar por paciente que habilita ensembles y fusión; control estricto de fuga."),
    ("Camino del proyecto (p1 → p27)",
     ["p11: Catálogo por paciente (56 features)",
      "p22: Calibración (Platt/Isotónica) y análisis por cohorte",
      "p24: LR elastic-net + Platt (p_img)",
      "p26/p26b: Intermodal LATE/MID + calibración por cohorte + política S2"],
     "Hitos que consolidan la base técnica hasta una decisión coste-sensata."),
    ("P24 — Modelo de imagen (LR-EN + Platt)",
     ["TEST: AUC 0.727 (ALL) · 0.754 (OAS1) · 0.750 (OAS2) · Brier 0.220 (ALL)",
      "Umbrales 5:1 (VAL): OAS1 0.435 · OAS2 0.332",
      "p_img como pilar de la fusión intermodal"],
     "P24 ofrece probabilidades calibradas y estables para la fusión posterior."),
    ("P26 — Intermodal: LATE vs MID",
     ["LATE = media(p_img, p_clin) → simple y robusto",
      "MID integra IMG56+clínico+p1 pero rinde peor por N/covariables",
      "TEST LATE: AUC 0.713 · PR-AUC 0.712 · Brier 0.234 (mejor que MID)"],
     "La simplicidad de LATE reduce riesgo de sobreajuste y facilita despliegue."),
    ("P26b/P27 — Calibración por cohorte + Política S2",
     ["Platt por cohorte y coste 5:1",
      "S2 ajusta OAS2 para Recall ≥ 0.90",
      "Umbrales activos: OAS1 0.42 · OAS2 ≈ 0.4929",
      "Smoke TEST: OAS1 R=0.70, OAS2 R=0.917; Costes 39 y 11"],
     "Mostramos capacidad de ajustar la política al objetivo clínico (sensibilidad alta)."),
    ("Demo: GUI (Streamlit) + API (FastAPI) + CLI",
     ["Streamlit: carga CSV, switch de política, sliders, métricas (TP/FP/TN/FN, Coste) y gráficos",
      "FastAPI: POST /predict con clinical+features o clinical+p_img",
      "CLI: compute_pimg_from_features.py · predict_end_to_end.py"],
     "Listo para demo en vivo: experiencia completa de extremo a extremo."),
    ("Impacto, límites y próximos pasos",
     ["Impacto: prioriza recall y reduce FN en cribado",
      "Límites: N reducido y ECE mayor en OAS2 → recalibración por sitio y monitoring",
      "Próximos: validación externa, domain adaptation, reporte clínico"],
     "Cierre con la historia humana: ganar tiempo de calidad para el paciente."),
    ("Backup — Comparativa de pipelines (TEST)",
     ["P19 ALL: AUC 0.671 · PR 0.606 · Brier 0.292",
      "P22 ALL: AUC 0.668–0.702 · PR 0.605–0.646",
      "P24 ALL: AUC 0.727 · PR 0.717 · Brier 0.220",
      "P26 LATE: AUC 0.713 · PR 0.712 · Brier 0.234"],
     "Refuerzo de la narrativa de mejora y maduración del sistema."),
    ("Backup — Decisión coste-sensible (5:1)",
     ["OAS1 @0.435: TP 14 · FP 9 · TN 18 · FN 6 (Coste 39)",
      "OAS2 @0.332: TP 11 · FP 7 · TN 4 · FN 1 (Coste 12)",
      "S2: OAS1 0.42 · OAS2 ≈ 0.4929 → OAS2 Recall ≥ 0.90"],
     "Tabla de confusión como soporte a decisiones clínicas."),
    ("Backup — Calibración (ECE/MCE)",
     ["TEST ECE@10/MCE: ALL 0.178/0.407",
      "OAS1 0.150/0.578 · OAS2 0.313/0.766",
      "Acción: Platt por cohorte + monitorizar drift (ECE>0.20 → recalibrar)"],
     "Demuestra atención a fiabilidad probabilística."),
    ("Backup — Arquitectura y release",
     ["Carpetas: models, CONFIG, QA, DOCS, MANIFEST/ENV",
      "Streamlit + FastAPI + CLI (end-to-end)",
      "Zip reproducible y portable"],
     "Cierre técnico: listo para evaluación y demo.")
]

for title, bullets, notes in slides_payload:
    add_slide(title, bullets, notes)

prs.save("CognitivaAI_Story_Guion.pptx")
print("✅ Generado: CognitivaAI_Story_Guion.pptx")